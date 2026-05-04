"""
Modify Slide 2 of the Innovate X NIET PPT:
- Update IDEA TITLE
- Update Proposed Solution text
- Resize Detailed Explanation box to left half
- Add compact bullet content with proper font size
- Right side left empty for flow diagram image (user adds manually or we add)
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import copy

prs = Presentation(r"c:\INNOVATE X HACKATHON\Innovate X NIET -IDEA-Presentation-Format.pptx")

slide = prs.slides[1]  # Slide 2 (0-indexed)

# Print all shapes to understand layout
print("=== SLIDE 2 SHAPES ===")
for i, shape in enumerate(slide.shapes):
    print(f"  [{i}] name='{shape.shape_id}:{shape.name}' | "
          f"left={shape.left}, top={shape.top}, "
          f"width={shape.width}, height={shape.height}")
    if shape.has_text_frame:
        for p in shape.text_frame.paragraphs:
            txt = p.text.strip()
            if txt:
                print(f"       TEXT: {txt[:80]}")

print("\n=== Modifying slide... ===")

# Find shapes by their text content
for shape in slide.shapes:
    if not shape.has_text_frame:
        continue
    full_text = " ".join(p.text for p in shape.text_frame.paragraphs).strip()

    # === UPDATE IDEA TITLE ===
    if "IDEA TITLE" in full_text:
        print(f"  Found IDEA TITLE shape: {shape.name}")
        for para in shape.text_frame.paragraphs:
            if "IDEA TITLE" in para.text:
                for run in para.runs:
                    run.text = ""
                para.runs[0].text = "TrustPath"
                # Keep existing formatting but make sure it's visible
                print("  -> Updated to 'TrustPath'")

    # === UPDATE PROPOSED SOLUTION ===
    if "Proposed Solution" in full_text and "Detailed" not in full_text:
        print(f"  Found Proposed Solution shape: {shape.name}")
        tf = shape.text_frame
        # Clear all paragraphs after the title
        # Find the paragraph with the description
        for para in tf.paragraphs:
            text = para.text.strip()
            if text and "Proposed Solution" not in text:
                for run in para.runs:
                    run.text = ""
                para.runs[0].text = (
                    "An AI-powered placement readiness platform that verifies "
                    "student skills through staged adaptive assessments and real "
                    "evidence, then creates personalized roadmaps, recruiter trust "
                    "profiles, and college-level analytics."
                )
                para.runs[0].font.size = Pt(14)
                print("  -> Updated proposed solution text")

    # === UPDATE DETAILED EXPLANATION ===
    if "Detailed explanation" in full_text or "detailed explanation" in full_text.lower():
        print(f"  Found Detailed Explanation shape: {shape.name}")
        print(f"    Current size: {shape.width}x{shape.height}")
        print(f"    Current pos: {shape.left},{shape.top}")

        # Resize to left 55% of usable width
        # Slide width = 12192000 EMU (10 inches standard)
        # Keep top position, adjust width
        shape.width = Emu(5800000)  # ~4.7 inches (left half)

        tf = shape.text_frame
        tf.word_wrap = True

        # Clear existing content
        while len(tf.paragraphs) > 0:
            p = tf.paragraphs[0]
            for run in p.runs:
                run.text = ""
            if len(tf.paragraphs) == 1:
                break
            # Remove extra paragraphs by clearing
            p._p.getparent().remove(p._p)

        # Bullet content - the sweet spot version
        bullets = [
            ("Skill Intake & Claim Extraction: ", "Resume upload (NLP-parsed for skills & projects) OR manual skill entry + target role (SDE, AI/ML) & dream company selection"),
            ("3-Stage Adaptive Test (IRT): ", "Easy (5m) → Medium (10m) → Hard (15m) — questions adapt in real-time; tracks accuracy, response time, answer changes & confidence self-ratings to detect skill bluffing"),
            ("ML Trust Engine (XGBoost): ", "7 signals — accuracy by difficulty, time patterns, confidence calibration, answer instability, Codeforces/LeetCode evidence, resume-claim alignment, project strength → Trust Score (0–100) + Bluff Risk Index"),
            ("Trust Stamp — Verified URL: ", "Consent-based shareable link on resume — recruiters see verified scores, skill evidence & risk flags. One click, full transparency. No fake claims."),
            ("Gated AI Roadmap: ", "Role-specific skill tree (node-based). Each node = curated resources + assignments + mini-projects. Must submit proof to unlock next — real skill acquisition only."),
            ("Gamification & Platform Sync: ", "XP, streaks, college leaderboards (opt-in). Syncs with LeetCode/Codeforces/GFG — level up based on actual problems solved."),
            ("Smart ATS Engine: ", "Assessment-fused resume scoring — \"ATS: 62%. Verified gaps: add React project, highlight CF rating (1400+).\""),
            ("Institutional Dashboard: ", "Real-time cohort analytics — skill gap heatmaps, bluff-risk distribution, placement readiness %. \"Why your batch fails at Company X.\""),
        ]

        # Set first paragraph (title line)
        first_p = tf.paragraphs[0]
        first_p.text = ""
        run = first_p.add_run()
        run.text = "Detailed Explanation"
        run.font.size = Pt(16)
        run.font.bold = True
        run.font.color.rgb = RGBColor(0, 0, 0)
        first_p.alignment = PP_ALIGN.LEFT

        # Add bullet paragraphs
        for bold_part, normal_part in bullets:
            p = tf.add_paragraph()
            p.alignment = PP_ALIGN.LEFT
            p.space_before = Pt(2)
            p.space_after = Pt(1)
            p.level = 0

            # Bold part
            r1 = p.add_run()
            r1.text = bold_part
            r1.font.size = Pt(9)
            r1.font.bold = True
            r1.font.color.rgb = RGBColor(0, 0, 0)

            # Normal part
            r2 = p.add_run()
            r2.text = normal_part
            r2.font.size = Pt(9)
            r2.font.bold = False
            r2.font.color.rgb = RGBColor(30, 30, 30)

        print(f"  -> Added {len(bullets)} bullet points at 9pt font")
        print(f"  -> New size: {shape.width}x{shape.height}")

# Save
output_path = r"c:\INNOVATE X HACKATHON\TrustPath_Presentation.pptx"
prs.save(output_path)
print(f"\n✅ Saved to: {output_path}")
